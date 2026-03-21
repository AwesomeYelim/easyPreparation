"""
슬라이드 생성 프레임워크
─────────────────────────────────────────────────────────────────────────
■ 이식 방법
  다른 프로젝트로 복사 후 _TEMPLATE, _LOGO 경로만 조정하면 끝.
  스타일 상수는 5개뿐이므로 필요 시 아래 블록만 수정.

■ 슬라이드 생성 패턴 (인라인 — 파일 저장 금지)
  tools/.venv/Scripts/python.exe - <<'EOF'   # Windows
  tools/.venv/bin/python - <<'EOF'           # Linux/macOS
  import sys; sys.path.insert(0, "tools/scripts")
  from generate_slide import *
  prs, LAYOUT_TITLE, LAYOUT_CONTENT = new_prs()
  s = add_slide(prs, LAYOUT_CONTENT)
  slide_title(s, "제목")
  tbl(s, ["컬럼1","컬럼2"], [("행1","값1")])
  logo(s)
  prs.save("tools/output/result.pptx")
  EOF

■ 배경 규칙
  LAYOUT_TITLE   → 레이아웃 배경 이미지 자동 적용 (add_slide 만으로 완성)
  LAYOUT_CONTENT → 흰색 배경 자동 적용
  bg_rect 추가 금지
"""

import os
import json
import tempfile
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE

# ── 공식 템플릿 경로 (절대 변경 금지) ────────────────────────────────────
_TEMPLATE = os.path.join(os.path.dirname(__file__), "..", "templates", "slide_template.pptx")
_LOGO     = os.path.join(os.path.dirname(__file__), "..", "templates", "ssrinc_logo.png")

# ── 스타일 상수 (5개) ─────────────────────────────────────────────────────
FONT    = "NanumGothic"  # 템플릿 폰트
C_WHITE = "FFFFFF"       # 빨간 배경 위 텍스트 (타이틀 슬라이드)
C_DARK  = "1A1A1A"       # 흰색 배경 위 텍스트 (콘텐츠 슬라이드)
C_RED   = "A31627"       # 브랜드 레드 — 섹션 제목·테이블 헤더·강조
C_GRAY  = "EEEEEE"       # 테이블 교차 행 배경
# ─────────────────────────────────────────────────────────────────────────

# 로고 배치 기준 (slide_template.pptx Slide 1 좌표 — 너비만 지정, 종횡비 자동)
_LOGO_X = Inches(8.22)
_LOGO_Y = Inches(6.42)
_LOGO_W = Inches(2.13)

# 슬라이드 크기 (템플릿에서 1회 읽어 모듈 상수로 고정)
_tmpl = Presentation(_TEMPLATE) if os.path.exists(_TEMPLATE) else None
SLIDE_W = _tmpl.slide_width  if _tmpl else Inches(10.83)
SLIDE_H = _tmpl.slide_height if _tmpl else Inches(7.5)
del _tmpl


# ── 프레젠테이션 생성 ─────────────────────────────────────────────────────

def new_prs():
    """템플릿 기반 Presentation 반환. 기존 슬라이드 제거, 마스터/레이아웃 유지.
    반환값: (prs, layout_title, layout_content)
    인라인 스크립트에서: prs, LAYOUT_TITLE, LAYOUT_CONTENT = new_prs()
    """
    if os.path.exists(_TEMPLATE):
        prs = Presentation(_TEMPLATE)
        sldIdLst = prs.slides._sldIdLst
        for sid in list(sldIdLst):
            prs.part.drop_rel(sid.rId)
            sldIdLst.remove(sid)
    else:
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H
    lt = prs.slide_layouts[0]
    lc = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]
    return prs, lt, lc


def add_slide(prs, layout):
    """슬라이드 추가 후 layout placeholder를 모두 제거하여 반환.
    prs.slides.add_slide() 직접 호출 금지 — 이 함수만 사용.
    placeholder 요소가 커스텀 도형 위를 덮어 내용이 안 보이는 문제를 방지.
    """
    slide = prs.slides.add_slide(layout)
    sp_tree = slide.shapes._spTree
    for ph in slide.placeholders:
        sp_tree.remove(ph._element)
    return slide


# ── 헬퍼 함수 (모듈 레벨 — import 가능) ──────────────────────────────────

def _rgb(h: str):
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def txt(slide, text: str, x, y, w, h,
        size: int = 16, bold: bool = False,
        color: str = C_DARK, align=PP_ALIGN.LEFT):
    """텍스트박스 추가. 폰트·색상은 상수 기본값 사용."""
    r, g, b = _rgb(color)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = RGBColor(r, g, b)
    return tb


def slide_title(slide, text: str):
    """콘텐츠 슬라이드 섹션 제목 — 브랜드 레드, 굵게."""
    txt(slide, text,
        Inches(0.37), Inches(0.12), SLIDE_W - Inches(0.74), Inches(0.6),
        size=26, bold=True, color=C_RED)


def logo(slide):
    """로고 삽입 — 종횡비 자동 유지 (width 만 지정)."""
    if os.path.exists(_LOGO):
        slide.shapes.add_picture(_LOGO, _LOGO_X, _LOGO_Y, width=_LOGO_W)


def tbl(slide, headers: list, rows: list,
        y_start=Inches(1.0), row_h=Inches(0.48), col_ws=None):
    """
    테이블 그리기.
    헤더 행: C_RED 배경 + 흰색 텍스트
    데이터 행: 흰색/연회색(C_GRAY) 교차
    col_ws: 각 열 너비 리스트 (미지정 시 균등 분할)
    """
    ncols = len(headers)
    if col_ws is None:
        col_ws = [(SLIDE_W - Inches(1.0)) / ncols] * ncols
    col_xs = [Inches(0.5)]
    for w in col_ws[:-1]:
        col_xs.append(col_xs[-1] + w)

    for ri, row in enumerate([headers] + list(rows)):
        y = y_start + ri * row_h
        is_hdr = (ri == 0)
        rbg = C_RED if is_hdr else (C_GRAY if ri % 2 == 1 else "FFFFFF")
        for ci, (val, cx, cw) in enumerate(zip(row, col_xs, col_ws)):
            cell = slide.shapes.add_shape(1, cx, y, cw, row_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*_rgb(rbg))
            cell.line.color.rgb = RGBColor(*_rgb("CCCCCC"))
            txt(slide, str(val),
                cx + Inches(0.05), y + Inches(0.05), cw - Inches(0.1), row_h - Inches(0.1),
                size=10 if not is_hdr else 11,
                bold=is_hdr,
                color=C_WHITE if is_hdr else C_DARK,
                align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)


def steps_list(slide, steps: list, y_start=Inches(1.0), row_h=Inches(0.72)):
    """
    번호 + 함수명 + 설명 형태의 순서 목록.
    steps: [(번호, 함수명, 설명), ...]
    """
    _CIRC = Inches(0.38)
    _TXT_H = Inches(0.42)
    for i, (num, name, desc) in enumerate(steps):
        y      = y_start + i * row_h
        cy     = y + (row_h - _CIRC)  / 2   # 원 수직 중앙 정렬
        ty     = y + (row_h - _TXT_H) / 2   # 텍스트 수직 중앙 정렬
        circ   = slide.shapes.add_shape(9, Inches(0.35), cy, _CIRC, _CIRC)
        circ.fill.solid()
        circ.fill.fore_color.rgb = RGBColor(*_rgb(C_RED))
        circ.line.width = Pt(0)
        txt(slide, str(num), Inches(0.35), cy, _CIRC, _CIRC,
            size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(slide, name, Inches(0.85), ty, Inches(4.3), _TXT_H,
            size=11, bold=True, color=C_RED)
        txt(slide, desc, Inches(5.3), ty, SLIDE_W - Inches(5.8), _TXT_H,
            size=11, color=C_DARK)


# ── 데모 build() — 실제 사용 시 인라인 스크립트로 대체 ────────────────────

def build(output_path: str, data: dict = None):
    """
    최소 데모 (2 슬라이드).
    실제 문서 생성은 이 함수를 쓰지 않고 인라인 스크립트에서
    new_prs() + 헬퍼 함수를 직접 호출한다.
    """
    prs, lt, lc = new_prs()

    # Slide 1: 타이틀 (레이아웃이 빨간 배경 자동 제공)
    s1 = add_slide(prs, lt)
    txt(s1, "문서 제목",
        Inches(0.5), Inches(2.9), SLIDE_W - Inches(1.0), Inches(1.0),
        size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s1, "부제목 · 날짜",
        Inches(0.5), Inches(4.2), SLIDE_W - Inches(1.0), Inches(0.6),
        size=14, color=C_WHITE, align=PP_ALIGN.CENTER)
    logo(s1)

    # Slide 2: 콘텐츠 예시 (레이아웃이 흰색 배경 자동 제공)
    s2 = add_slide(prs, lc)
    slide_title(s2, "섹션 제목")
    tbl(s2,
        ["항목", "값"],
        [("예시 항목 1", "값 1"),
         ("예시 항목 2", "값 2")],
    )
    logo(s2)

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("output", nargs="?", default=None)
    parser.add_argument("--data", default=None)
    args = parser.parse_args()

    data = {}
    if args.data and os.path.exists(args.data):
        with open(args.data, encoding="utf-8") as f:
            data = json.load(f)

    out = args.output or os.path.join(tempfile.gettempdir(), "output.pptx")
    print(build(out, data))


if __name__ == "__main__":
    main()
